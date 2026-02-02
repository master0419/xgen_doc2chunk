"""
PPT Chart Extractor

Extracts chart data from PowerPoint files (PPTX).
Uses python-pptx Presentation and Chart objects.

Provides:
- extract(): Single chart extraction from python-pptx Chart object
- extract_all_from_file(): Extract all charts from PPTX file in slide order
"""
import io
import logging
from typing import Any, BinaryIO, Dict, List, Optional, Union

from pptx import Presentation

from xgen_doc2chunk.core.functions.chart_extractor import BaseChartExtractor, ChartData

logger = logging.getLogger("document-processor")


class PPTChartExtractor(BaseChartExtractor):
    """
    Chart extractor for PowerPoint files.
    
    Supports:
    - Direct python-pptx Chart object extraction
    - Full file extraction via extract_all_from_file()
    """
    
    # ========================================================================
    # Main Interface
    # ========================================================================
    
    def extract(self, chart_element: Any) -> ChartData:
        """
        Extract chart data from python-pptx Chart object.
        
        Args:
            chart_element: python-pptx Chart object
                
        Returns:
            ChartData with extracted information
        """
        if not chart_element:
            return ChartData()
        
        title = self._extract_title(chart_element)
        chart_type = self._extract_chart_type(chart_element)
        categories = self._extract_categories(chart_element)
        series = self._extract_series(chart_element)
        
        return ChartData(
            chart_type=chart_type,
            title=title,
            categories=categories,
            series=series
        )
    
    def extract_all_from_file(
        self,
        file_source: Union[str, bytes, BinaryIO]
    ) -> List[ChartData]:
        """
        Extract all charts from a PowerPoint file in slide order.
        
        Args:
            file_source: File path, bytes, or file-like object
            
        Returns:
            List of ChartData for all charts in the file
        """
        charts = []
        
        try:
            # Prepare file-like object
            if isinstance(file_source, str):
                with open(file_source, 'rb') as f:
                    file_obj = io.BytesIO(f.read())
            elif isinstance(file_source, bytes):
                file_obj = io.BytesIO(file_source)
            else:
                file_source.seek(0)
                file_obj = file_source
            
            # Open presentation
            prs = Presentation(file_obj)
            
            # Iterate slides in order
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_chart:
                        chart_data = self.extract(shape.chart)
                        charts.append(chart_data)
                    
                    # Check group shapes
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, 'has_chart') and sub_shape.has_chart:
                                chart_data = self.extract(sub_shape.chart)
                                charts.append(chart_data)
            
            logger.info(f"Extracted {len(charts)} charts from PowerPoint file")
            
        except Exception as e:
            logger.error(f"Error extracting charts from PowerPoint: {e}")
        
        return charts
    
    # ========================================================================
    # Private Methods
    # ========================================================================
    
    def _extract_title(self, chart) -> Optional[str]:
        """Extract chart title."""
        try:
            if chart.has_title and chart.chart_title:
                if chart.chart_title.has_text_frame:
                    title_text = chart.chart_title.text_frame.text
                    if title_text:
                        return title_text.strip()
        except Exception:
            pass
        return None
    
    def _extract_chart_type(self, chart) -> str:
        """Extract chart type."""
        try:
            if hasattr(chart, 'chart_type'):
                type_str = str(chart.chart_type)
                type_name = type_str.split('.')[-1].split(' ')[0]
                return type_name.replace('_', ' ').title()
        except Exception:
            pass
        return "Chart"
    
    def _extract_categories(self, chart) -> List[str]:
        """Extract category labels."""
        categories = []
        try:
            if hasattr(chart, 'plots') and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, 'categories') and plot.categories:
                        categories = [str(c) for c in plot.categories]
                        break
        except Exception:
            pass
        return categories
    
    def _extract_series(self, chart) -> List[Dict[str, Any]]:
        """Extract series data."""
        series_data = []
        try:
            for idx, series in enumerate(chart.series):
                series_info = {
                    'name': self._get_series_name(series, idx),
                    'values': []
                }
                
                try:
                    if hasattr(series, 'values') and series.values:
                        series_info['values'] = list(series.values)
                except Exception:
                    pass
                
                series_data.append(series_info)
        except Exception:
            pass
        
        return series_data
    
    def _get_series_name(self, series, idx: int) -> str:
        """Get series name."""
        try:
            if hasattr(series, 'name') and series.name:
                return str(series.name)
        except Exception:
            pass
        return f"Series {idx + 1}"


__all__ = ['PPTChartExtractor']
