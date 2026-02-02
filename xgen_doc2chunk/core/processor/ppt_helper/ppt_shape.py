"""
PPT Shape 처리 모듈

포함 함수:
- get_shape_position(): Shape의 위치 정보 반환
- is_picture_shape(): Shape이 이미지인지 확인
- process_image_shape(): 이미지 Shape 처리 및 로컬 저장
- process_group_shape(): 그룹 Shape 처리
"""
import logging
from typing import Any, Dict, List, Optional, Tuple

from xgen_doc2chunk.core.functions.img_processor import ImageProcessor

from xgen_doc2chunk.core.processor.ppt_helper.ppt_constants import ElementType, SlideElement
from xgen_doc2chunk.core.processor.ppt_helper.ppt_bullet import extract_text_with_bullets
from xgen_doc2chunk.core.processor.ppt_helper.ppt_table import is_simple_table, extract_simple_table_as_text, convert_table_to_html

logger = logging.getLogger("document-processor")


def get_shape_position(shape) -> Tuple[int, int, int, int]:
    """
    Shape의 위치 정보를 반환합니다.

    Args:
        shape: python-pptx Shape 객체

    Returns:
        (left, top, width, height) 튜플 (EMU 단위)
    """
    try:
        left = shape.left if hasattr(shape, 'left') and shape.left else 0
        top = shape.top if hasattr(shape, 'top') and shape.top else 0
        width = shape.width if hasattr(shape, 'width') and shape.width else 0
        height = shape.height if hasattr(shape, 'height') and shape.height else 0
        return (left, top, width, height)
    except Exception:
        return (0, 0, 0, 0)


def is_picture_shape(shape) -> bool:
    """
    Shape이 이미지인지 확인합니다.

    Args:
        shape: python-pptx Shape 객체

    Returns:
        이미지이면 True
    """
    # 방법 1: shape_type 확인
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass

    # 방법 2: image 속성 확인
    if hasattr(shape, 'image'):
        try:
            _ = shape.image
            return True
        except Exception:
            pass

    return False


def process_image_shape(
    shape,
    processed_images: set,
    image_processor: ImageProcessor
) -> Optional[str]:
    """
    이미지 Shape을 처리하고 로컬에 저장합니다.

    Args:
        shape: python-pptx Shape 객체 (이미지)
        processed_images: 이미 처리된 이미지 해시 집합
        image_processor: ImageProcessor 인스턴스

    Returns:
        이미지 태그 문자열 또는 None
    """

    try:
        if not hasattr(shape, 'image'):
            return None

        image = shape.image
        image_bytes = image.blob

        if not image_bytes:
            return None

        image_tag = image_processor.save_image(image_bytes, processed_images=processed_images)

        if image_tag:
            return f"\n{image_tag}\n"

        return None

    except Exception as e:
        logger.warning("Error processing image shape: %s", e)
        return None


def process_group_shape(
    group_shape,
    processed_images: set,
    image_processor: ImageProcessor
) -> List[SlideElement]:
    """
    그룹 Shape 내의 요소들을 처리합니다.

    Args:
        group_shape: python-pptx Group Shape 객체
        processed_images: 이미 처리된 이미지 해시 집합
        image_processor: ImageProcessor 인스턴스

    Returns:
        SlideElement 리스트
    """

    elements = []

    try:
        for shape in group_shape.shapes:
            position = get_shape_position(shape)
            shape_id = shape.shape_id if hasattr(shape, 'shape_id') else id(shape)

            if shape.has_table:
                # 단순 표(1xN, Nx1, 2x2 이하)는 텍스트로 처리
                if is_simple_table(shape.table):
                    simple_text = extract_simple_table_as_text(shape.table)
                    if simple_text:
                        elements.append(SlideElement(
                            element_type=ElementType.TEXT,
                            content=simple_text,
                            position=position,
                            shape_id=shape_id
                        ))
                else:
                    # 일반 표는 HTML로 처리
                    table_html = convert_table_to_html(shape.table)
                    if table_html:
                        elements.append(SlideElement(
                            element_type=ElementType.TABLE,
                            content=table_html,
                            position=position,
                            shape_id=shape_id
                        ))

            elif is_picture_shape(shape):
                image_tag = process_image_shape(shape, processed_images, image_processor)
                if image_tag:
                    elements.append(SlideElement(
                        element_type=ElementType.IMAGE,
                        content=image_tag,
                        position=position,
                        shape_id=shape_id
                    ))

            # 텍스트 처리 - 목록 정보 포함
            elif hasattr(shape, "text_frame") and shape.text_frame:
                text_content = extract_text_with_bullets(shape.text_frame)
                if text_content:
                    elements.append(SlideElement(
                        element_type=ElementType.TEXT,
                        content=text_content,
                        position=position,
                        shape_id=shape_id
                    ))

            # 기존 text 속성만 있는 경우 (폴백)
            elif hasattr(shape, "text") and shape.text.strip():
                elements.append(SlideElement(
                    element_type=ElementType.TEXT,
                    content=shape.text.strip(),
                    position=position,
                    shape_id=shape_id
                ))

    except Exception as e:
        logger.warning(f"Error processing group shape: {e}")

    return elements
